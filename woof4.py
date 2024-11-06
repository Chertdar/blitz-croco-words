import os
import zipfile
from typing import IO
from pptx import Presentation

from pathlib import Path

FILE_TXT: str = 'words.txt'
FILE_ZIP = 'src/croco-blitz-source.zip'


def get_words_presentation(file: IO[bytes]) -> list[str]:
    prs = Presentation(file)
    words_list = list()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if ' ' in shape.text or ':' in shape.text or '-' in shape.text:
                continue
            words_list.append(shape.text)

    return words_list


def write_txt(filename: str, words: list[str]):
    try:
        with open(filename, 'w', encoding='utf-8') as file_text:
            for word in words:
                file_text.write(word + '\n')
        print(f"Слова записаны в файл: {filename}")
    except IOError as e:
        print(f"error: {e}")


def unzip_file():
    cur_folder = os.path.dirname(os.path.realpath(__file__))

    words: set[str] = set()

    p = Path(cur_folder) / FILE_ZIP

    with zipfile.ZipFile(p) as zip_ref:
        for f in zip_ref.namelist():
            words.update(get_words_presentation(zip_ref.open(f)))


if __name__ == '__main__':
    write_txt(unzip_file(), FILE_TXT)
