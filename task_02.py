# Поменяйте код так, чтобы он выводил текст каждого шэйпа на экран.
# Сохраните из этой презентации только слова, всякие "Правила", "Тур 1" и другие объяснение не нужны.
# Сохраните все слова в файл words.txt
# Закомитьте его тоже.
import re

from pptx import Presentation

from pathlib import Path

FILE_TXT: str = 'words.txt'


def get_words_presentation():
    prs = Presentation("Osennyaya_igra_11.pptx")
    words_list = list()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if ' ' in shape.text or ':' in shape.text or '-' in shape.text:
                continue
            words_list.append(shape.text)

    return words_list


def write_txt(filename: str, words: list[str]):
    try:
        with open(filename, 'w', encoding='utf-8') as file_text:
            for word in words:
                file_text.write(word + '\n')
        print(f"Слова записаны в файл: {filename}")
    except IOError as e:
        print(f"error: {e}")


if __name__ == '__main__':
    write_txt(FILE_TXT, get_words_presentation())
